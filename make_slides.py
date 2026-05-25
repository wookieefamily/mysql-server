#!/usr/bin/env python3
"""Generate analyst-vs-AI investigation workflow slides."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

C_ANALYST = RGBColor(0x3D, 0x3D, 0x3D)   # dark charcoal
C_AI      = RGBColor(0xEB, 0x1C, 0x24)   # Adobe red
C_CARD    = RGBColor(0xF0, 0xF0, 0xF0)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK    = RGBColor(0x1A, 0x1A, 0x1A)
C_GRAY    = RGBColor(0x66, 0x66, 0x66)
C_ARROW   = RGBColor(0xAA, 0xAA, 0xAA)
C_TRACK   = RGBColor(0xCC, 0xCC, 0xCC)
C_BORDER  = RGBColor(0xCC, 0xCC, 0xCC)
C_ACCENT  = RGBColor(0x99, 0x99, 0x99)   # top card accent bar

# MSO auto-shape type IDs
RECT   = 1
RDRECT = 5
OVAL   = 9
RARROW = 13

STEPS = [
    ("1", ["Identify a", "Trend Change"],
     "Watch dashboards. Monitor KPIs. Catch what's worth investigating."),
    ("2", ["Setup Report", "& Investigate"],
     "Find segments. Pull data. Build the chart. Spot the obvious driver."),
    ("3", ["Find", "the Why"],
     "Run breakdowns. Slack partner teams. Search Jira. Read briefs. Wait for replies."),
    ("4", ["Recommend", "& Write Up"],
     "Synthesize. Write it up. Get review. Send to leader."),
]

SLIDES = [
    dict(
        title="What analysis & investigation work looks like today",
        subtitle="Every step is performed manually by the analyst.",
        splits=[(1.0, 0.0), (1.0, 0.0), (1.0, 0.0), (1.0, 0.0)],
    ),
    dict(
        title="How AI transforms the investigation workflow",
        subtitle="Steps 1 & 2 fully automated by AI. Steps 3 & 4 AI-assisted.",
        splits=[(0.0, 1.0), (0.0, 1.0), (0.10, 0.90), (0.25, 0.75)],
    ),
]


def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def set_para(tf, lines, size, bold=False, color=C_DARK, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color


def draw_meter(slide, x, y, w, h, apct, aipct):
    x, y, w, h = int(x), int(y), int(w), int(h)

    # Grey track background
    bg = slide.shapes.add_shape(RECT, x, y, w, h)
    solid(bg, C_TRACK)
    no_line(bg)

    # Analyst segment (left, charcoal)
    if apct > 0.005:
        aw = max(1, int(w * apct))
        s = slide.shapes.add_shape(RECT, x, y, aw, h)
        solid(s, C_ANALYST)
        no_line(s)
        if apct > 0.14:
            tf = s.text_frame
            tf.margin_top = Pt(2)
            tf.margin_bottom = Pt(0)
            tf.margin_left = Pt(2)
            tf.margin_right = Pt(2)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = f"{int(round(apct * 100))}%"
            r.font.size = Pt(8)
            r.font.bold = True
            r.font.color.rgb = C_WHITE

    # AI segment (right, red)
    if aipct > 0.005:
        off = int(w * apct)
        aiw = w - off
        s = slide.shapes.add_shape(RECT, x + off, y, max(1, aiw), h)
        solid(s, C_AI)
        no_line(s)
        if aipct > 0.14:
            tf = s.text_frame
            tf.margin_top = Pt(2)
            tf.margin_bottom = Pt(0)
            tf.margin_left = Pt(2)
            tf.margin_right = Pt(2)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = f"{int(round(aipct * 100))}%"
            r.font.size = Pt(8)
            r.font.bold = True
            r.font.color.rgb = C_WHITE


def make_slide(prs, cfg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # White background
    bg = slide.shapes.add_shape(RECT, 0, 0, SLIDE_W, SLIDE_H)
    solid(bg, C_WHITE)
    no_line(bg)

    # Red left-edge accent bar (like the original slide)
    accent = slide.shapes.add_shape(RECT, 0, Inches(0.15), Inches(0.07), SLIDE_H - Inches(0.3))
    solid(accent, C_AI)
    no_line(accent)

    # Slide title
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(12.4), Inches(0.6))
    set_para(tb.text_frame, [cfg["title"]], 22, bold=True)

    # Subtitle
    sb = slide.shapes.add_textbox(Inches(0.45), Inches(0.82), Inches(12.4), Inches(0.35))
    set_para(sb.text_frame, [cfg["subtitle"]], 11, color=C_GRAY)

    # Card layout
    CW  = Inches(2.85)
    CH  = Inches(5.6)
    CY  = Inches(1.25)
    GAP = Inches(0.28)
    N   = 4
    START = (SLIDE_W - (N * CW + (N - 1) * GAP)) / 2

    for i, (num, title_lines, desc) in enumerate(STEPS):
        apct, aipct = cfg["splits"][i]
        cx = START + i * (CW + GAP)

        # Card
        card = slide.shapes.add_shape(RDRECT, int(cx), int(CY), int(CW), int(CH))
        solid(card, C_CARD)
        card.line.color.rgb = C_BORDER
        card.line.width = Pt(0.5)

        # Top accent bar on card (thin colored stripe at top)
        top_color = C_AI if aipct > apct else C_ANALYST
        tacc = slide.shapes.add_shape(RECT, int(cx), int(CY), int(CW), int(Inches(0.07)))
        solid(tacc, top_color)
        no_line(tacc)

        # Step number circle
        CSZ = Inches(0.65)
        cx_ = cx + (CW - CSZ) / 2
        cy_ = CY + Inches(0.18)
        circ = slide.shapes.add_shape(OVAL, int(cx_), int(cy_), int(CSZ), int(CSZ))
        solid(circ, C_DARK)
        no_line(circ)
        tf = circ.text_frame
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = C_WHITE

        # Step title
        ty = cy_ + CSZ + Inches(0.14)
        ttb = slide.shapes.add_textbox(
            int(cx + Inches(0.1)), int(ty),
            int(CW - Inches(0.2)), int(Inches(0.9))
        )
        set_para(ttb.text_frame, title_lines, 13, bold=True, align=PP_ALIGN.CENTER)

        # Description
        dy = ty + Inches(0.92)
        dtb = slide.shapes.add_textbox(
            int(cx + Inches(0.2)), int(dy),
            int(CW - Inches(0.4)), int(Inches(2.0))
        )
        set_para(dtb.text_frame, [desc], 10, color=C_GRAY, align=PP_ALIGN.CENTER)

        # ---- Meter section ----
        mpad = Inches(0.2)
        mx   = cx + mpad
        mw   = CW - 2 * mpad
        mh   = Inches(0.34)
        ly   = CY + CH - Inches(0.90)   # label row y
        my   = ly + Inches(0.26)        # meter bar y

        # "Analyst" label (left, charcoal)
        lb = slide.shapes.add_textbox(int(mx), int(ly), int(mw * 0.6), int(Inches(0.24)))
        tf = lb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "Analyst"
        r.font.size = Pt(8)
        r.font.bold = True
        r.font.color.rgb = C_ANALYST

        # "AI" label (right, red)
        rb = slide.shapes.add_textbox(
            int(mx + mw * 0.4), int(ly),
            int(mw * 0.6), int(Inches(0.24))
        )
        tf = rb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = "AI"
        r.font.size = Pt(8)
        r.font.bold = True
        r.font.color.rgb = C_AI

        # Meter bar
        draw_meter(slide, mx, my, mw, mh, apct, aipct)

        # Arrow to next card
        if i < N - 1:
            arr = slide.shapes.add_shape(
                RARROW,
                int(cx + CW + Inches(0.02)),
                int(CY + CH / 2 - Inches(0.15)),
                int(GAP - Inches(0.04)),
                int(Inches(0.3))
            )
            solid(arr, C_ARROW)
            no_line(arr)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for cfg in SLIDES:
        make_slide(prs, cfg)

    out = "/home/user/mysql-server/analysis_workflow.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
